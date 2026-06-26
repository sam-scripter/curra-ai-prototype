"""
parser.py — Document text extraction service.

This is the first step in the ingestion pipeline. It reads a raw file
and returns the text content as a structured list of pages.

Each supported format needs its own extraction approach:
  - PDF:  PyMuPDF reads each page's text layer directly
  - PPTX: python-pptx iterates over slides and their text frames
  - DOCX: python-docx reads paragraphs and groups them into logical pages

The output of this module — a list of ParsedPage dicts — is the input
to the chunker. The interface between them is deliberately simple so
either side can be improved independently.
"""

import os
import fitz                      # PyMuPDF — imported as fitz (its original name)
from pptx import Presentation    # python-pptx
from docx import Document        # python-docx
from typing import TypedDict


class ParsedPage(TypedDict):
    """
    One page or slide worth of extracted text.

    page_number: 1-indexed position in the document.
                 PDF  → actual page number
                 PPTX → slide number
                 DOCX → logical group number (see parse_docx)
    text:        Raw extracted text. May contain newlines and extra whitespace.
                 The chunker handles cleaning.
    """
    page_number: int
    text: str


def parse_pdf(file_path: str) -> list[ParsedPage]:
    """
    Extract text from a PDF page by page using PyMuPDF.

    PyMuPDF is preferred over pypdf for lecture slides because it handles
    complex layouts, embedded fonts, and scanned-but-OCR'd PDFs better.
    get_text("text") returns plain text — other modes like "html" or "dict"
    give richer structure but more noise for our use case.
    """
    pages: list[ParsedPage] = []
    doc = fitz.open(file_path)

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()

        # Skip blank pages — common in slide decks (title slides, dividers).
        # A blank page produces an empty string and would create a useless chunk.
        if text:
            pages.append({"page_number": page_num, "text": text})

    doc.close()
    return pages


def parse_pptx(file_path: str) -> list[ParsedPage]:
    """
    Extract text from a PowerPoint file slide by slide.

    Each slide is treated as one page. We extract from two sources:
    1. Text frames — the visible content boxes and titles on the slide
    2. Speaker notes — the notes panel below the slide, which often contains
       the lecturer's explanations. These are labelled [Lecturer Notes]
       so the LLM can distinguish them from slide content.

    Extracting notes significantly increases the value of each slide
    for RAG because lecturers often explain concepts in notes that are
    only bullet-pointed on the slide itself.
    """
    pages: list[ParsedPage] = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []

        # Extract text from every shape on the slide that contains text.
        # Shapes include: title placeholder, content placeholder, text boxes,
        # tables. We only handle text frames here — images are out of scope
        # for the prototype (Phase 2 of the full system adds GPT-4o Vision).
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)

        # Extract speaker notes if the slide has them.
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                if notes_text:
                    parts.append(f"[Lecturer Notes] {notes_text}")

        combined = "\n".join(parts).strip()
        if combined:
            pages.append({"page_number": slide_num, "text": combined})

    return pages


def parse_docx(file_path: str) -> list[ParsedPage]:
    """
    Extract text from a Word document.

    DOCX files have no inherent page concept — pages are determined by
    rendering, not by the file structure. We group paragraphs into logical
    batches of 30 so that the page_number in citations stays meaningful
    (rather than all chunks citing "page 1").
    """
    PARAGRAPHS_PER_PAGE = 30

    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    pages: list[ParsedPage] = []

    for i in range(0, len(paragraphs), PARAGRAPHS_PER_PAGE):
        batch = paragraphs[i:i + PARAGRAPHS_PER_PAGE]
        page_num = (i // PARAGRAPHS_PER_PAGE) + 1
        pages.append({
            "page_number": page_num,
            "text": "\n".join(batch),
        })

    return pages


def parse_file(file_path: str) -> list[ParsedPage]:
    """
    Main entry point for the parser.

    Detects the file format from the extension and delegates to the
    appropriate parser function. Raises ValueError for unsupported formats
    so the caller (ingestion pipeline) can handle it cleanly.
    """
    extension = os.path.splitext(file_path)[1].lower()

    parsers = {
        ".pdf":  parse_pdf,
        ".pptx": parse_pptx,
        ".docx": parse_docx,
    }

    if extension not in parsers:
        raise ValueError(
            f"Unsupported file format '{extension}'. "
            f"Supported: {', '.join(parsers.keys())}"
        )

    return parsers[extension](file_path)